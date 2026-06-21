"""Generates a polished hackathon pitch deck for DeepTraffic-Guard.
Run: python3 build_pitch_deck.py
Output: DeepTraffic-Guard_Pitch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0B, 0x0F, 0x1A)
SURFACE = RGBColor(0x11, 0x18, 0x27)
SURFACE2 = RGBColor(0x1A, 0x22, 0x35)
BORDER = RGBColor(0x1E, 0x2D, 0x47)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SUBTLE = RGBColor(0x47, 0x55, 0x69)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT2 = RGBColor(0x60, 0xA5, 0xFA)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
YELLOW = RGBColor(0xEA, 0xB3, 0x08)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)

IMG_DIR = "/home/anmol/.claude/image-cache/c8670732-bdf8-4066-a1c0-526e451341d5"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide(bg_color=NAVY):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def no_autofit(tf):
    try:
        tf.word_wrap = True
        tf.auto_size = None
    except Exception:
        pass


def add_text(slide, x, y, w, h, text, size=18, color=TEXT, bold=False,
             italic=False, align=PP_ALIGN.LEFT, font="Calibri", anchor=None,
             line_spacing=1.0, letter_spacing=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    no_autofit(tf)
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, size=15, color=TEXT, accent=ACCENT,
                 gap=Pt(10), marker="▸"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    no_autofit(tf)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        if isinstance(item, tuple):
            label, desc = item
            r1 = p.add_run()
            r1.text = f"{marker} " + label + "  "
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = accent
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(size - 1)
            r2.font.bold = False
            r2.font.color.rgb = MUTED
            r2.font.name = "Calibri"
        else:
            r1 = p.add_run()
            r1.text = f"{marker} " + item
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            r1.font.name = "Calibri"
    return box


def add_card(slide, x, y, w, h, fill=SURFACE, line=BORDER, radius=0.06, line_w=1.0):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(line_w)
    card.shadow.inherit = False
    try:
        card.adjustments[0] = radius
    except Exception:
        pass
    return card


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_dot(slide, x, y, d, color):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    dot.shadow.inherit = False
    return dot


def accent_bar(slide, x, y, w=Inches(0.55), h=Pt(4), color=ACCENT):
    return add_rect(slide, x, y, w, h, color)


def kicker_pill(slide, x, y, text, color=ACCENT):
    w = Inches(0.18 * len(text) + 0.5)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.34))
    pill.fill.solid()
    pill.fill.fore_color.rgb = SURFACE2
    pill.line.color.rgb = color
    pill.line.width = Pt(1)
    pill.shadow.inherit = False
    try:
        pill.adjustments[0] = 0.5
    except Exception:
        pass
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = color
    return pill


def slide_title(slide, kicker, title, x=Inches(0.7), y=Inches(0.45), title_size=31):
    kicker_pill(slide, x, y, kicker)
    add_text(slide, x, y + Inches(0.5), Inches(11.8), Inches(0.8), title,
              size=title_size, color=TEXT, bold=True)
    accent_bar(slide, x, y + Inches(1.18))


def page_chrome(slide, n, section=""):
    add_text(slide, Inches(0.7), Inches(7.08), Inches(4), Inches(0.3),
              "DEEPTRAFFIC-GUARD", size=9, color=SUBTLE, bold=True)
    if section:
        add_text(slide, Inches(5.5), Inches(7.08), Inches(2.3), Inches(0.3),
                  section, size=9, color=SUBTLE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(12.4), Inches(7.08), Inches(0.6), Inches(0.3),
              f"{n:02d}", size=9, color=SUBTLE, align=PP_ALIGN.RIGHT)


def corner_accent(slide):
    add_rect(slide, 0, 0, Inches(0.12), SH, ACCENT)


# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
s = add_slide()
corner_accent(s)
add_dot(s, Inches(6.27), Inches(1.55), Inches(0.8), SURFACE2)
add_text(s, Inches(0), Inches(1.6), Inches(13.333), Inches(0.7), "🚦", size=34,
          align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(2.75), Inches(13.333), Inches(1.0), "DeepTraffic-Guard",
          size=58, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.75), Inches(13.333), Inches(0.55),
          "Turning Every Traffic Camera Into a 24/7 Enforcement Officer",
          size=20, color=ACCENT2, align=PP_ALIGN.CENTER, bold=True)
add_text(s, Inches(0), Inches(4.35), Inches(13.333), Inches(0.45),
          "Real-Time AI Violation Detection  •  Zero Manual Setup  •  Evidence-Grade Output",
          size=13.5, color=MUTED, align=PP_ALIGN.CENTER)

tags = ["YOLOv8", "OpenCV", "EasyOCR", "Streamlit", "8 Violations"]
tag_gap = Inches(0.22)
tag_widths = [Inches(0.18 * len(t) + 0.5) for t in tags]  # must match kicker_pill's formula
total = sum(tag_widths, Inches(0)) + tag_gap * (len(tags) - 1)
tx = (SW - total) / 2
for tag, w in zip(tags, tag_widths):
    kicker_pill(s, tx, Inches(5.15), tag, color=ACCENT2)
    tx += w + tag_gap
add_text(s, Inches(0), Inches(6.7), Inches(13.333), Inches(0.4),
          "Live Demo  →  deeptraffic-guard.streamlit.app", size=13,
          color=MUTED, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 2 — THE HOOK
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "The Reality", "Right now, somewhere in your city...", title_size=34)
hook_items = [
    "A rider with no helmet just passed an intersection camera.",
    "Three people are crammed onto one motorcycle.",
    "A car just ran a red light at full speed.",
]
y = Inches(2.3)
for i, item in enumerate(hook_items):
    add_card(s, Inches(0.9), y, Inches(11.3), Inches(0.95), fill=SURFACE)
    add_dot(s, Inches(1.15), y + Inches(0.3), Inches(0.35), [RED, ORANGE, RED][i])
    add_text(s, Inches(1.75), y + Inches(0.23), Inches(10.2), Inches(0.5), item,
              size=17, color=TEXT)
    y += Inches(1.1)
add_text(s, Inches(0.9), Inches(5.85), Inches(11.3), Inches(0.9),
          "Nobody saw any of it. The camera recorded it — but no one was watching.",
          size=19, color=ACCENT2, bold=True, italic=True)
page_chrome(s, 2, "The Problem")

# ============================================================================
# SLIDE 3 — THE PROBLEM, QUANTIFIED
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "The Problem", "Enforcement capacity hasn't scaled with traffic")
add_bullets(s, Inches(0.7), Inches(2.05), Inches(6.3), Inches(4.3), [
    ("One officer, one junction.", "The moment they look away, or move to another junction, violations elsewhere go completely unnoticed."),
    ("The cameras already exist.", "Most cities have traffic cameras on nearly every major junction — footage is recorded, but rarely actively monitored in real time."),
    ("The stakes are lives, not fines.", "Helmet, seatbelt, and red-light laws exist because ignoring them is a leading cause of fatal accidents — not because of revenue."),
], size=15.5, gap=Pt(14))

card = add_card(s, Inches(7.35), Inches(2.05), Inches(5.25), Inches(4.4), fill=SURFACE2)
add_text(s, Inches(7.65), Inches(2.3), Inches(4.7), Inches(0.4), "THE GAP",
          size=13, color=ACCENT2, bold=True)
add_text(s, Inches(7.65), Inches(2.8), Inches(4.7), Inches(1.2),
          "Footage exists.\nWatching doesn't.", size=27, color=TEXT, bold=True,
          line_spacing=1.25)
add_rect(s, Inches(7.65), Inches(4.15), Inches(4.4), Pt(1.2), BORDER)
add_text(s, Inches(7.65), Inches(4.4), Inches(4.4), Inches(1.9),
          "A system that watches every frame, never gets tired, never looks away — "
          "and turns footage that already exists into evidence that gets acted on.",
          size=13.5, color=MUTED, line_spacing=1.35)
page_chrome(s, 3, "The Problem")

# ============================================================================
# SLIDE 4 — SOLUTION OVERVIEW
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Our Solution", "One photo in. Full evidence out.")
add_text(s, Inches(0.7), Inches(2.05), Inches(11.8), Inches(0.85),
          "Upload a traffic camera photo or video. DeepTraffic-Guard automatically finds "
          "every vehicle and person, checks 8 violation types, reads license plates, and "
          "produces evidence-grade output — with zero manual setup per camera.",
          size=16, color=MUTED, line_spacing=1.3)

steps = [("📷", "Upload"), ("🧹", "Auto-Enhance"), ("🎯", "Detect Objects"),
         ("🛑", "Find Road Rules"), ("⚖️", "Check Violations"), ("📄", "Evidence + PDF")]
x0 = Inches(0.7)
gap = Inches(2.0)
for i, (icon, label) in enumerate(steps):
    x = x0 + i * gap
    add_card(s, x, Inches(3.25), Inches(1.78), Inches(1.65), fill=SURFACE)
    add_text(s, x, Inches(3.4), Inches(1.78), Inches(0.6), icon, size=26, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(4.1), Inches(1.78), Inches(0.6), label, size=11.5,
              color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(4.45), Inches(1.78), Inches(0.3), f"STEP {i+1}", size=8,
              color=SUBTLE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s, x + Inches(1.8), Inches(3.65), Inches(0.22), Inches(0.8), "→",
                  size=20, color=ACCENT, align=PP_ALIGN.CENTER)

add_card(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.35), fill=SURFACE2)
add_text(s, Inches(1.0), Inches(5.45), Inches(0.5), Inches(1.0), "★", size=24,
          color=ACCENT2)
add_text(s, Inches(1.5), Inches(5.46), Inches(10.8), Inches(1.05),
          "Key differentiator: stop lines, no-parking zones, and road signs are detected "
          "directly from the photo's road markings — never assumed, never hand-configured. "
          "The same system deploys instantly at a brand-new camera.",
          size=14.5, color=ACCENT2, bold=True, line_spacing=1.25)
page_chrome(s, 4, "Solution")

# ============================================================================
# SLIDE 5 — VIOLATIONS GRID
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Core Detection", "8 violation types, one unified pipeline")
violations = [
    ("Helmet Non-Compliance", "Skin / hair / texture HSV analysis on rider head crop", RED),
    ("Triple Riding", "Rider count linked per motorcycle bounding box", RED),
    ("Red-Light Violation", "Traffic-light color check + stop-line crossing", RED),
    ("Wrong-Side Driving", "Lane-center estimation via Hough line transform", RED),
    ("Phone Use While Driving", "Cell-phone bbox overlap with driver / rider", RED),
    ("Stop-Line Violation", "Vehicle position vs. auto-detected stop line", ORANGE),
    ("Seatbelt Non-Compliance", "Diagonal-strap edge density in torso region", ORANGE),
    ("Illegal Parking", "Vehicle overlap with auto-detected curb-paint zones", YELLOW),
]
cols = 2
cw, ch = Inches(5.85), Inches(0.95)
gx, gy = Inches(0.25), Inches(0.15)
ox, oy = Inches(0.7), Inches(2.0)
for i, (label, desc, color) in enumerate(violations):
    r, c = divmod(i, cols)
    x = ox + c * (cw + gx)
    y = oy + r * (ch + gy)
    add_card(s, x, y, cw, ch, fill=SURFACE)
    add_dot(s, x + Inches(0.22), y + Inches(0.32), Inches(0.3), color)
    add_text(s, x + Inches(0.68), y + Inches(0.09), cw - Inches(0.85), Inches(0.4), label,
              size=13.5, color=TEXT, bold=True)
    add_text(s, x + Inches(0.68), y + Inches(0.44), cw - Inches(0.85), Inches(0.45), desc,
              size=10, color=MUTED)
page_chrome(s, 5, "Detection")

# ============================================================================
# SLIDE 6 — DEMO (INPUT)
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Live Demo · Input", "A real, unedited Bengaluru traffic photo")
try:
    s.shapes.add_picture(f"{IMG_DIR}/7.png", Inches(0.7), Inches(2.0), height=Inches(4.95))
except Exception:
    pass
add_card(s, Inches(8.3), Inches(2.0), Inches(4.3), Inches(4.95), fill=SURFACE2)
add_text(s, Inches(8.6), Inches(2.25), Inches(3.8), Inches(0.4), "ZERO SETUP",
          size=12, color=ACCENT2, bold=True)
add_text(s, Inches(8.6), Inches(2.7), Inches(3.8), Inches(1.0),
          "A location the system\nhas never seen before.",
          size=16, color=TEXT, bold=True, line_spacing=1.3)
add_bullets(s, Inches(8.6), Inches(3.75), Inches(3.8), Inches(1.5), [
    "No stop line marked",
    "No parking zone marked",
    "No lane config",
], size=13.5, color=TEXT, marker="—")
add_rect(s, Inches(8.6), Inches(5.35), Inches(3.5), Pt(1), BORDER)
add_text(s, Inches(8.6), Inches(5.55), Inches(3.8), Inches(1.3),
          "12 objects in frame:\nbuses, cars, motorcycles,\nriders, auto-rickshaws",
          size=12.5, color=MUTED, line_spacing=1.3)
page_chrome(s, 6, "Live Demo")

# ============================================================================
# SLIDE 7 — DEMO (RESULT)
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Live Demo · Output", "Same photo, processed in 7 milliseconds")
try:
    s.shapes.add_picture(f"{IMG_DIR}/8.png", Inches(0.7), Inches(2.0), height=Inches(4.95))
except Exception:
    pass
page_chrome(s, 7, "Live Demo")

# ============================================================================
# SLIDE 8 — ENGINEERING RIGOR / DEBUG STORY
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Engineering Rigor", "We don't just demo it — we stress-tested it")
add_text(s, Inches(0.7), Inches(2.0), Inches(11.8), Inches(0.6),
          "Real photos break naive computer vision. Here's a bug we found and fixed by "
          "testing against actual traffic photos, not synthetic data:",
          size=14.5, color=MUTED, line_spacing=1.3)

add_card(s, Inches(0.7), Inches(2.75), Inches(5.75), Inches(3.9), fill=SURFACE)
add_text(s, Inches(1.0), Inches(2.95), Inches(5.2), Inches(0.4), "🐛 THE BUG", size=13,
          color=RED, bold=True)
add_text(s, Inches(1.0), Inches(3.4), Inches(5.2), Inches(3.1),
          "A rider wearing a bright red helmet was being flagged as 'no helmet'.\n\n"
          "Root cause: pure red and fair skin tones share the same hue in HSV color "
          "space. Glossy highlights on the helmet shell were misread as exposed skin.\n\n"
          "Deeper still: our own shadow-correction preprocessing step was silently "
          "desaturating colors across the whole image — masking the fix until we "
          "traced it end-to-end.",
          size=12.5, color=MUTED, line_spacing=1.35)

add_card(s, Inches(6.7), Inches(2.75), Inches(5.95), Inches(3.9), fill=SURFACE2)
add_text(s, Inches(7.0), Inches(2.95), Inches(5.4), Inches(0.4), "✅ THE FIX", size=13,
          color=GREEN, bold=True)
add_text(s, Inches(7.0), Inches(3.4), Inches(5.4), Inches(3.1),
          "Added a 'helmet-shell rejection' signal: large saturated-color regions "
          "are recognized as painted shells, not skin.\n\n"
          "Fixed shadow correction to adjust brightness only (HSV value channel), "
          "preserving true color for every downstream check.\n\n"
          "Verified against the original photo + 13 other test images, plus "
          "synthetic boundary cases — zero regressions.",
          size=12.5, color=MUTED, line_spacing=1.35)
page_chrome(s, 8, "Engineering")

# ============================================================================
# SLIDE 9 — ARCHITECTURE
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Architecture", "Eight decoupled, independently testable modules")
stages = [
    ("Image /\nVideo Input", SURFACE2), ("Preprocessor", SURFACE), ("YOLOv8\nDetector", SURFACE2),
    ("Rider\nAssociation", SURFACE), ("Road Zone\nAuto-Detect", SURFACE2), ("Violation\nEngine", SURFACE),
    ("OCR +\nAnnotator", SURFACE2), ("PDF /\nDashboard", SURFACE),
]
n = len(stages)
total_w = Inches(12.0)
box_w = Inches(1.32)
arrow_w = (total_w - box_w * n) / (n - 1)
x = Inches(0.65)
y = Inches(2.9)
for i, (label, color) in enumerate(stages):
    add_card(s, x, y, box_w, Inches(1.35), fill=color)
    add_text(s, x, y + Inches(0.18), box_w, Inches(1.0), label, size=10.5,
              color=TEXT, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < n - 1:
        add_text(s, x + box_w, y + Inches(0.38), arrow_w, Inches(0.6), "→",
                  size=17, color=ACCENT, align=PP_ALIGN.CENTER)
    x += box_w + arrow_w

add_card(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.7), fill=SURFACE2)
add_text(s, Inches(1.0), Inches(5.05), Inches(11.3), Inches(1.4),
          "Each module is a plain Python file with a single responsibility — preprocessing, "
          "detection, zone-finding, and violation logic never touch each other's internals. "
          "That separation is what let us isolate, reproduce, and fix the helmet bug in "
          "minutes instead of guessing across a monolith.",
          size=14, color=MUTED, line_spacing=1.35)
page_chrome(s, 9, "Architecture")

# ============================================================================
# SLIDE 10 — TECH STACK
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Tech Stack", "Built entirely on proven, open-source tools")
stack = [
    ("YOLOv8", "Real-time object detection", "Ultralytics"),
    ("OpenCV", "HSV, Canny, Hough, CLAHE", "Classic CV"),
    ("EasyOCR", "License plate text recognition", "Deep learning OCR"),
    ("Streamlit", "Interactive web UI, zero frontend code", "Python-native"),
    ("SQLite", "Violation storage + analytics", "Embedded DB"),
    ("fpdf2", "Auto-generated PDF evidence reports", "Reporting"),
]
cols = 3
cw, ch = Inches(3.85), Inches(1.75)
gx, gy = Inches(0.25), Inches(0.3)
ox, oy = Inches(0.7), Inches(2.1)
for i, (label, desc, tag) in enumerate(stack):
    r, c = divmod(i, cols)
    x = ox + c * (cw + gx)
    y = oy + r * (ch + gy)
    add_card(s, x, y, cw, ch, fill=SURFACE)
    add_text(s, x + Inches(0.25), y + Inches(0.18), cw - Inches(0.5), Inches(0.45), label,
              size=17, color=ACCENT2, bold=True)
    add_text(s, x + Inches(0.25), y + Inches(0.68), cw - Inches(0.5), Inches(0.65), desc,
              size=11, color=MUTED, line_spacing=1.2)
    kicker_pill(s, x + Inches(0.25), y + Inches(1.3), tag, color=SUBTLE)
page_chrome(s, 10, "Tech Stack")

# ============================================================================
# SLIDE 11 — COMPETITIVE EDGE
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Why We're Different", "Most CV hackathon projects stop at detection")
rows = [
    ("Setup per new camera", "Manual zone calibration", "Auto-detected from road paint"),
    ("Output format", "Raw bounding boxes", "Annotated evidence + PDF report"),
    ("Input type", "Single image only", "Image AND sampled video"),
    ("Violation coverage", "Usually 1-2 types", "8 independent violation checks"),
    ("Testing rigor", "Demo on curated images", "Stress-tested on real, messy photos"),
]
headers = ["", "Typical Hackathon Project", "DeepTraffic-Guard"]
col_x = [Inches(0.7), Inches(4.0), Inches(8.65)]
col_w = [Inches(3.1), Inches(4.4), Inches(4.4)]
y = Inches(2.05)
for i, h in enumerate(headers):
    color = TEXT if i == 0 else (MUTED if i == 1 else GREEN)
    add_text(s, col_x[i], y, col_w[i], Inches(0.4), h, size=12.5, color=color, bold=True)
y += Inches(0.5)
for label, bad, good in rows:
    add_card(s, Inches(0.7), y, Inches(12.0), Inches(0.78), fill=SURFACE)
    add_text(s, col_x[0] + Inches(0.2), y + Inches(0.16), col_w[0], Inches(0.5), label,
              size=12.5, color=TEXT, bold=True)
    add_text(s, col_x[1] + Inches(0.2), y + Inches(0.16), col_w[1] - Inches(0.2), Inches(0.5),
              "✗  " + bad, size=12, color=MUTED)
    add_text(s, col_x[2] + Inches(0.2), y + Inches(0.16), col_w[2] - Inches(0.2), Inches(0.5),
              "✓  " + good, size=12, color=GREEN, bold=True)
    y += Inches(0.9)
page_chrome(s, 11, "Differentiation")

# ============================================================================
# SLIDE 12 — BEYOND DETECTION
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Beyond Detection", "A complete, submission-ready system")
add_bullets(s, Inches(0.7), Inches(2.1), Inches(11.8), Inches(4.4), [
    ("Video analysis mode", "Frame-sampled scanning of full traffic video clips, with an annotated output video and a per-type violation timeline — not just single photos."),
    ("One-click PDF evidence reports", "Professional report generation: annotated frame, severity-coded violation table, and a dedicated evidence page per violation."),
    ("Full analytics dashboard", "Violations by type/severity/hour, searchable record history, CSV export, and automatic frame-hash deduplication."),
    ("Zero manual configuration", "Stop lines, parking zones, and signs are detected per-photo from actual road paint — the same system deploys instantly at any new camera."),
], size=15.5, gap=Pt(15))
page_chrome(s, 12, "Completeness")

# ============================================================================
# SLIDE 13 — IMPACT
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Impact & Scalability", "Existing infrastructure, new capability")
metrics = [
    ("24/7", "Coverage per camera —\nno shift changes, no fatigue"),
    ("8", "Independent violation\ntypes checked per frame"),
    ("0", "Manual setup steps\nrequired per new camera"),
]
cw = Inches(3.8)
gx = Inches(0.3)
ox = Inches(0.85)
for i, (num, label) in enumerate(metrics):
    x = ox + i * (cw + gx)
    add_card(s, x, Inches(2.25), cw, Inches(2.5), fill=SURFACE2)
    add_text(s, x, Inches(2.5), cw, Inches(1.0), num, size=50, color=ACCENT2,
              bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.55), cw, Inches(1.0), label, size=13.5, color=MUTED,
              align=PP_ALIGN.CENTER, line_spacing=1.3)
add_card(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(1.55), fill=SURFACE)
add_text(s, Inches(1.0), Inches(5.35), Inches(11.3), Inches(1.2),
          "Cities don't need new hardware to deploy this — every existing traffic camera "
          "becomes an active enforcement layer. Footage that already exists turns into "
          "evidence that already gets used.",
          size=15, color=TEXT, line_spacing=1.35)
page_chrome(s, 13, "Impact")

# ============================================================================
# SLIDE 14 — LIMITATIONS / ROADMAP
# ============================================================================
s = add_slide()
corner_accent(s)
slide_title(s, "Honest Limitations & Roadmap", "What's next")
add_text(s, Inches(0.7), Inches(2.0), Inches(5.7), Inches(0.4), "TODAY'S LIMITATIONS",
          size=13.5, color=RED, bold=True)
add_bullets(s, Inches(0.7), Inches(2.5), Inches(5.7), Inches(4.0), [
    ("Resolution floor", "Tiny/distant riders have too few pixels for reliable head analysis."),
    ("Heuristic vs. trained models", "Helmet/seatbelt use rule-based CV rather than a dedicated trained classifier — by design, for transparency."),
    ("Multi-lane highways", "Wrong-side check assumes a simple two-lane road model."),
], size=14, gap=Pt(13))

add_text(s, Inches(6.9), Inches(2.0), Inches(5.7), Inches(0.4), "ROADMAP",
          size=13.5, color=GREEN, bold=True)
add_bullets(s, Inches(6.9), Inches(2.5), Inches(5.7), Inches(4.0), [
    ("Trained helmet classifier", "Fine-tune a model on labeled helmet/no-helmet photos for higher accuracy at distance."),
    ("Live camera feed ingestion", "Connect directly to RTSP/live streams instead of file upload."),
    ("Multi-camera fleet dashboard", "Aggregate analytics across many camera nodes, city-wide."),
], size=14, accent=GREEN, gap=Pt(13))
page_chrome(s, 14, "Roadmap")

# ============================================================================
# SLIDE 15 — THANK YOU / CTA
# ============================================================================
s = add_slide()
corner_accent(s)
add_text(s, Inches(0), Inches(2.3), Inches(13.333), Inches(1.0), "See It Live",
          size=48, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.25), Inches(13.333), Inches(0.5),
          "DeepTraffic-Guard — AI Traffic Enforcement System", size=17,
          color=ACCENT2, align=PP_ALIGN.CENTER)
add_card(s, Inches(4.4), Inches(4.1), Inches(4.5), Inches(0.85), fill=SURFACE2,
         line=ACCENT)
add_text(s, Inches(4.4), Inches(4.3), Inches(4.5), Inches(0.5),
          "deeptraffic-guard.streamlit.app", size=16, color=ACCENT2, bold=True,
          align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.15), Inches(13.333), Inches(0.4),
          "github.com/anmol-2-4/deeptraffic-guard", size=13, color=MUTED,
          align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.3), Inches(13.333), Inches(0.5),
          "Thank you — questions welcome.", size=15, color=MUTED, italic=True,
          align=PP_ALIGN.CENTER)

prs.save("DeepTraffic-Guard_Pitch.pptx")
print("done")
